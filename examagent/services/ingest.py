"""Document loading and chunking.

Supported: PDF, TXT, MD, DOCX, PPTX, and images (registered but not text-
extracted unless an OCR engine is installed). Every chunk keeps the metadata
needed for an honest citation: source_type, source_name, lecture, page, section.
"""
from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

from ..config import get_logger, get_settings
from ..data.topics import TOPIC_SEEDS

log = get_logger(__name__)

TEXT_EXT = {".txt", ".md", ".markdown", ".rst", ".tex", ".csv"}
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
SUPPORTED = TEXT_EXT | PDF_EXT | DOCX_EXT | PPTX_EXT | IMAGE_EXT


@dataclass
class RawSection:
    """A contiguous piece of a document with its location metadata."""

    text: str
    page: int | None = None
    section: str | None = None


@dataclass
class Chunk:
    chunk_id: str
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


# --------------------------------------------------------------- loaders
def _load_text(path: Path) -> list[RawSection]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    sections: list[RawSection] = []
    current_head: str | None = None
    buf: list[str] = []
    for line in raw.splitlines():
        head = re.match(r"^\s{0,3}(#{1,4})\s+(.*)$", line)
        if head:
            if buf:
                sections.append(RawSection("\n".join(buf).strip(), None, current_head))
                buf = []
            current_head = head.group(2).strip()
            continue
        buf.append(line)
    if buf:
        sections.append(RawSection("\n".join(buf).strip(), None, current_head))
    return [s for s in sections if s.text.strip()] or [RawSection(raw, None, None)]


def _load_pdf(path: Path) -> list[RawSection]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    out: list[RawSection] = []
    for i, page in enumerate(reader.pages, 1):
        try:
            txt = page.extract_text() or ""
        except Exception as exc:  # pragma: no cover - malformed pdfs
            log.warning("pdf page %s failed: %s", i, exc)
            txt = ""
        if txt.strip():
            out.append(RawSection(txt, page=i))
    return out


def _load_docx(path: Path) -> list[RawSection]:
    import docx

    doc = docx.Document(str(path))
    out: list[RawSection] = []
    current_head: str | None = None
    buf: list[str] = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        if (p.style.name or "").lower().startswith("heading"):
            if buf:
                out.append(RawSection("\n".join(buf), None, current_head))
                buf = []
            current_head = text
            continue
        buf.append(text)
    if buf:
        out.append(RawSection("\n".join(buf), None, current_head))
    for table in doc.tables:
        rows = [" | ".join(c.text.strip() for c in r.cells) for r in table.rows]
        if any(r.strip(" |") for r in rows):
            out.append(RawSection("\n".join(rows), None, "table"))
    return out


def _load_pptx(path: Path) -> list[RawSection]:
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[RawSection] = []
    for i, slide in enumerate(prs.slides, 1):
        title = None
        parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title is None and shape == getattr(slide.shapes, "title", None):
                title = text
            parts.append(text)
        if parts:
            out.append(RawSection("\n".join(parts), page=i, section=title or f"Slide {i}"))
    return out


def _load_image(path: Path) -> list[RawSection]:
    """Best-effort OCR; returns an empty list when no OCR engine is available."""
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore

        text = pytesseract.image_to_string(Image.open(str(path)))
        return [RawSection(text, section="OCR")] if text.strip() else []
    except Exception as exc:
        log.info("no OCR for %s (%s)", path.name, type(exc).__name__)
        return []


def load_document(path: Path) -> list[RawSection]:
    ext = path.suffix.lower()
    if ext in PDF_EXT:
        return _load_pdf(path)
    if ext in DOCX_EXT:
        return _load_docx(path)
    if ext in PPTX_EXT:
        return _load_pptx(path)
    if ext in IMAGE_EXT:
        return _load_image(path)
    if ext in TEXT_EXT:
        return _load_text(path)
    raise ValueError(f"Unsupported file type: {ext}")


# --------------------------------------------------------------- chunking
_WS = re.compile(r"[ \t]+")
_NL = re.compile(r"\n{3,}")


def clean_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\x00", "")
    text = _WS.sub(" ", text)
    text = _NL.sub("\n\n", text)
    # de-hyphenate line breaks common in PDF extraction
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    return text.strip()


def split_paragraphs(text: str) -> list[str]:
    parts = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    return parts or ([text.strip()] if text.strip() else [])


def chunk_sections(
    sections: Iterable[RawSection],
    base_metadata: dict[str, Any],
    chunk_size: int | None = None,
    overlap: int | None = None,
) -> list[Chunk]:
    """Pack paragraphs into ~chunk_size character chunks with overlap.

    Chunks never span a page/section boundary, so page citations stay truthful.
    """
    settings = get_settings()
    size = chunk_size or settings.chunk_size
    ov = overlap if overlap is not None else settings.chunk_overlap

    chunks: list[Chunk] = []
    for sec in sections:
        text = clean_text(sec.text)
        if len(text) < 40:
            continue
        paras = split_paragraphs(text)
        buf: list[str] = []
        buf_len = 0

        def flush() -> None:
            nonlocal buf, buf_len
            if not buf:
                return
            body = "\n\n".join(buf).strip()
            if len(body) >= 40:
                meta = dict(base_metadata)
                meta["page"] = sec.page
                meta["section"] = sec.section
                cid = hashlib.sha1(
                    (str(meta.get("source_name")) + str(sec.page) + body[:200]
                     + str(len(chunks))).encode()
                ).hexdigest()[:20]
                chunks.append(Chunk(cid, body, meta))
            buf = []
            buf_len = 0

        for para in paras:
            # a single oversized paragraph is hard-split on sentence boundaries
            if len(para) > size * 1.6:
                flush()
                sentences = re.split(r"(?<=[.!?])\s+", para)
                cur: list[str] = []
                cur_len = 0
                for s in sentences:
                    if cur_len + len(s) > size and cur:
                        buf = [" ".join(cur)]
                        buf_len = cur_len
                        flush()
                        tail = " ".join(cur)[-ov:] if ov else ""
                        cur = [tail, s] if tail else [s]
                        cur_len = len(tail) + len(s)
                    else:
                        cur.append(s)
                        cur_len += len(s) + 1
                if cur:
                    buf = [" ".join(cur)]
                    buf_len = cur_len
                    flush()
                continue

            if buf_len + len(para) > size and buf:
                tail = "\n\n".join(buf)[-ov:] if ov else ""
                flush()
                if tail:
                    buf = [tail]
                    buf_len = len(tail)
            buf.append(para)
            buf_len += len(para) + 2
        flush()

    return chunks


# --------------------------------------------------------------- topic tagging
_TOPIC_PATTERNS: list[tuple[str, list[str]]] = []
for _t in TOPIC_SEEDS:
    terms = [_t["name"].lower()] + [k.lower() for k in _t.get("keywords", [])]
    _TOPIC_PATTERNS.append((_t["id"], terms))


def detect_topics(text: str, limit: int = 6) -> list[str]:
    """Cheap keyword tagging so retrieval and coverage stats can filter by topic."""
    low = text.lower()
    scored: list[tuple[str, int]] = []
    for tid, terms in _TOPIC_PATTERNS:
        hits = sum(low.count(term) for term in terms if len(term) > 3)
        if hits:
            scored.append((tid, hits))
    scored.sort(key=lambda kv: -kv[1])
    return [t for t, _ in scored[:limit]]


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for block in iter(lambda: fh.read(1 << 16), b""):
            h.update(block)
    return h.hexdigest()


def prepare_document(
    path: Path,
    source_type: str,
    source_name: str | None = None,
    lecture: str = "",
) -> tuple[list[Chunk], dict[str, Any]]:
    """Load + chunk a file, returning chunks and a summary of what was found."""
    sections = load_document(path)
    n_chars = sum(len(s.text) for s in sections)
    base_meta = {
        "source_type": source_type,
        "source_name": source_name or path.stem,
        "lecture": lecture,
        "filename": path.name,
    }
    chunks = chunk_sections(sections, base_meta)
    joined = "\n".join(c.text for c in chunks[:60])
    topics = detect_topics(joined)
    for c in chunks:
        c.metadata["topics"] = ",".join(detect_topics(c.text, limit=4))
    summary = {
        "n_sections": len(sections),
        "n_chunks": len(chunks),
        "n_chars": n_chars,
        "topics": topics,
        "hash": file_hash(path),
    }
    return chunks, summary
